#!/usr/bin/env python3
import argparse
import re
from pathlib import Path
from typing import Dict, List, Tuple, Optional

from pptx import Presentation
from pptx.util import Pt


def read_markdown(md_path: Path) -> List[str]:
	md_text = md_path.read_text(encoding="utf-8")
	return md_text.splitlines()


def parse_sections(lines: List[str]) -> Tuple[str, Dict[str, List[str]]]:
	"""
	Parses a markdown file with:
	- H2 as document title (## Title)
	- H3 as section headers (### Section)
	- Bullets start with "- "
	Returns (doc_title, sections_dict)
	"""
	title = ""
	sections: Dict[str, List[str]] = {}
	current_section: Optional[str] = None

	for raw in lines:
		line = raw.rstrip()
		# Title (first H2)
		if line.startswith("## "):
			if not title:
				title = line[3:].strip()
			# do not treat H2 as a section
			continue
		# Sections (H3)
		if line.startswith("### "):
			current_section = line[4:].strip()
			sections[current_section] = []
			continue
		# Accumulate content lines under current section
		if current_section:
			sections[current_section].append(line)

	return title, sections


def clean_section_lines(lines: List[str]) -> List[str]:
	"""
	Cleans leading/trailing empty lines and returns meaningful content.
	"""
	# Trim leading/trailing blanks
	start = 0
	while start < len(lines) and lines[start].strip() == "":
		start += 1
	end = len(lines) - 1
	while end >= 0 and lines[end].strip() == "":
		end -= 1
	if end < start:
		return []
	return lines[start : end + 1]


def split_bullets_and_paragraphs(lines: List[str]) -> Tuple[List[str], List[str]]:
	"""
	Splits given lines into bullet items (starting with "- ") and paragraph text.
	- Bullet lines have their leading "- " removed
	- Paragraphs are joined into a single paragraph (preserving inline text)
	"""
	bullets: List[str] = []
	paragraph_lines: List[str] = []
	for line in lines:
		if line.strip().startswith("- "):
			item = re.sub(r"^\s*-\s+", "", line).strip()
			if item:
				bullets.append(item)
		else:
			if line.strip() != "":
				paragraph_lines.append(line.strip())
	return bullets, paragraph_lines


def add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None) -> None:
	slide_layout = prs.slide_layouts[0]  # Title slide
	slide = prs.slides.add_slide(slide_layout)
	slide.shapes.title.text = title
	if subtitle is not None and subtitle.strip():
		slide.placeholders[1].text = subtitle
		# Slightly larger subtitle
		for p in slide.placeholders[1].text_frame.paragraphs:
			for run in p.runs:
				run.font.size = Pt(22)


def add_content_slide(prs: Presentation, heading: str, bullets: List[str], paragraph: Optional[str]) -> None:
	slide_layout = prs.slide_layouts[1]  # Title and Content
	slide = prs.slides.add_slide(slide_layout)
	slide.shapes.title.text = heading
	text_frame = slide.shapes.placeholders[1].text_frame
	text_frame.clear()

	def add_bullet(text: str, level: int = 0) -> None:
		p = text_frame.add_paragraph() if len(text_frame.paragraphs) > 0 else text_frame.paragraphs[0]
		p.text = text
		p.level = level
		for run in p.runs:
			run.font.size = Pt(20)

	if paragraph:
		add_bullet(paragraph, level=0)
	# Add an empty line between paragraph and bullets when both exist
	if paragraph and bullets:
		add_bullet("", level=0)
	for b in bullets:
		add_bullet(b, level=0)


def guess_subtitle_from_title(doc_title: str) -> str:
	# Try to extract a month/year if present
	m = re.search(r"\b([A-Za-z]+)\s+20\d{2}\b", doc_title)
	if m:
		return f"{m.group(0)}"
	# Fallback generic subtitle
	return "Incident Management Monthly Report"


def build_presentation(input_md: Path, output_pptx: Path) -> None:
	lines = read_markdown(input_md)
	doc_title, sections = parse_sections(lines)
	prs = Presentation()

	# Title slide
	add_title_slide(prs, title=doc_title or "Incident Management Report", subtitle=guess_subtitle_from_title(doc_title))

	# Section slides
	for section_name, raw_lines in sections.items():
		content_lines = clean_section_lines(raw_lines)
		bullets, paragraph_lines = split_bullets_and_paragraphs(content_lines)
		paragraph_text = " ".join(paragraph_lines).strip() if paragraph_lines else None
		add_content_slide(prs, heading=section_name, bullets=bullets, paragraph=paragraph_text)

	# Ensure parent dir exists
	output_pptx.parent.mkdir(parents=True, exist_ok=True)
	prs.save(str(output_pptx))


def main():
	parser = argparse.ArgumentParser(description="Generate PowerPoint from Incident Management markdown report.")
	parser.add_argument("--input", "-i", type=Path, required=True, help="Path to input markdown file")
	parser.add_argument("--output", "-o", type=Path, required=True, help="Path to output .pptx file")
	args = parser.parse_args()

	build_presentation(args.input, args.output)


if __name__ == "__main__":
	main()


